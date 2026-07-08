from __future__ import annotations

import logging
from io import BytesIO
from typing import Callable

from .schemas import AttachmentMetadata, CustomAttachmentParser, FileType, _MEDIA_TYPE_MAP

logger = logging.getLogger("HarborRAG Connector::Confluence")


def _classify(media_type: str, title: str) -> tuple[FileType, str] | None:
    """Resolve (FileType, extension), with filename-based overrides for the
    media types Confluence reports ambiguously.
    """
    lowered_title = title.lower()

    if (
        media_type
        in (
            "application/vnd.ms-excel",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "text/plain",
        )
        and lowered_title.endswith(".csv")
    ):
        return FileType.CSV, "csv"

    if lowered_title.endswith((".md", ".mdx")) and media_type in (
        "binary/octet-stream",
        "text/markdown",
        "text/plain",
    ):
        return FileType.MARKDOWN, "md"

    return _MEDIA_TYPE_MAP.get(media_type)


class AttachmentProcessor:
    """Downloads and extracts text from a Confluence page's attachments.

    Framework-free: no BaseReader/Document/instrumentation-dispatcher
    dependency. Each format's extractor lazily imports its own dependencies
    so installing this package doesn't force every optional dependency
    (pytesseract, pandas, python-pptx, ...) on someone who only needs a
    subset of attachment types.
    """

    def __init__(
        self,
        *,
        download_fn: Callable[[str], bytes | None],
        base_url: str,
        custom_parsers: dict[FileType, CustomAttachmentParser] | None = None,
        process_attachment_callback: Callable[[str, int, str], tuple[bool, str]]
        | None = None,
        max_attachment_size_bytes: int | None = None,
        fail_on_error: bool = False,
        logger_: logging.Logger | None = None,
    ) -> None:
        self._download = download_fn
        self.base_url = base_url.rstrip("/")
        self.custom_parsers = custom_parsers or {}
        self.process_attachment_callback = process_attachment_callback
        self.max_attachment_size_bytes = max_attachment_size_bytes
        self.fail_on_error = fail_on_error
        self.logger = logger_ or logger

    def process(self, attachments: list[dict]) -> list[AttachmentMetadata]:
        """`attachments` is the raw `results` list from
        `client.get_attachments_from_content(page_id)`.
        """
        return [self._process_one(attachment) for attachment in attachments]

    def _process_one(self, attachment: dict) -> AttachmentMetadata:
        att_id = attachment.get("id", "")
        title = attachment.get("title", "")
        media_type = attachment.get("metadata", {}).get("mediaType", "")
        size_bytes = attachment.get("extensions", {}).get("fileSize", 0)
        download_path = attachment.get("_links", {}).get("download", "")
        download_url = f"{self.base_url}{download_path}"

        base = AttachmentMetadata(
            id=att_id,
            title=title,
            media_type=media_type,
            size_bytes=size_bytes,
            download_url=download_url,
            status="skipped",
        )

        if self.process_attachment_callback:
            should_process, reason = self.process_attachment_callback(
                media_type, size_bytes, title
            )
            if not should_process:
                base.reason = reason
                return base

        if self.max_attachment_size_bytes and size_bytes > self.max_attachment_size_bytes:
            base.reason = (
                f"size {size_bytes}B exceeds max_attachment_size_bytes "
                f"({self.max_attachment_size_bytes}B)"
            )
            return base

        classified = _classify(media_type, title)
        if classified is None:
            base.status = "unsupported"
            base.reason = f"no handler for media_type {media_type!r}"
            return base

        file_type, extension = classified

        try:
            content = self._download(download_url)
            if content is None:
                base.status = "failed"
                base.reason = "download failed or returned no content"
                return base

            text = (
                self.custom_parsers[file_type](content, extension)
                if file_type in self.custom_parsers
                else self._extract(file_type, content, extension)
            )
            base.status = "processed"
            base.text = text
            return base
        except Exception as exc:  # noqa: BLE001 - one bad attachment shouldn't sink the page
            self.logger.error("Failed to process attachment %s (%s): %s", title, att_id, exc)
            if self.fail_on_error:
                raise
            base.status = "failed"
            base.reason = str(exc)
            return base

    def _extract(self, file_type: FileType, content: bytes, extension: str) -> str:
        if file_type is FileType.SPREADSHEET:
            return self._extract_spreadsheet(content, extension)

        handler = {
            FileType.PDF: self._extract_pdf,
            FileType.IMAGE: self._extract_image,
            FileType.SVG: self._extract_svg,
            FileType.DOCUMENT: self._extract_docx,
            FileType.CSV: self._extract_csv,
            FileType.MESSAGE: self._extract_msg,
            FileType.HTML: self._extract_html,
            FileType.TEXT: self._extract_text,
            FileType.MARKDOWN: self._extract_text,
            FileType.PRESENTATION: self._extract_pptx,
        }[file_type]
        return handler(content)

    def _extract_pdf(self, content: bytes) -> str:
        try:
            import pytesseract
            from pdf2image import convert_from_bytes
        except ImportError as exc:
            raise ImportError(
                "`pytesseract` and `pdf2image` are required for PDF attachments; "
                "run `pip install pytesseract pdf2image`"
            ) from exc

        images = convert_from_bytes(content)
        pages = [
            f"Page {i + 1}:\n{pytesseract.image_to_string(image)}"
            for i, image in enumerate(images)
        ]
        return "\n\n".join(pages)

    def _extract_image(self, content: bytes) -> str:
        try:
            import pytesseract
            from PIL import Image
        except ImportError as exc:
            raise ImportError(
                "`pytesseract` and `Pillow` are required for image attachments; "
                "run `pip install pytesseract Pillow`"
            ) from exc

        return pytesseract.image_to_string(Image.open(BytesIO(content)))

    def _extract_svg(self, content: bytes) -> str:
        try:
            import pytesseract
            from PIL import Image
            from reportlab.graphics import renderPM
            from svglib.svglib import svg2rlg
        except ImportError as exc:
            raise ImportError(
                "`pytesseract`, `Pillow`, `svglib`, and `reportlab` are required for "
                "SVG attachments; run `pip install pytesseract Pillow svglib reportlab`"
            ) from exc

        drawing = svg2rlg(BytesIO(content))
        img_buffer = BytesIO()
        renderPM.drawToFile(drawing, img_buffer, fmt="PNG")
        img_buffer.seek(0)
        return pytesseract.image_to_string(Image.open(img_buffer))

    def _extract_docx(self, content: bytes) -> str:
        try:
            import docx2txt
        except ImportError as exc:
            raise ImportError(
                "`docx2txt` is required for Word attachments; run `pip install docx2txt`"
            ) from exc

        return docx2txt.process(BytesIO(content))

    def _extract_pptx(self, content: bytes) -> str:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ImportError(
                "`python-pptx` is required for PowerPoint attachments; "
                "run `pip install python-pptx`"
            ) from exc

        presentation = Presentation(BytesIO(content))
        chunks = [
            shape.text
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        ]
        return " ".join(chunks).strip()

    _EXCEL_ENGINES = {"xls": "xlrd", "xlsx": "openpyxl", "xlsm": "openpyxl", "xlsb": "pyxlsb"}

    def _extract_spreadsheet(self, content: bytes, extension: str) -> str:
        try:
            import pandas as pd
        except ImportError as exc:
            raise ImportError(
                "`pandas` is required for spreadsheet attachments; " \
                "run `pip install pandas`"
            ) from exc

        engine = self._EXCEL_ENGINES.get(extension)
        try:
            sheets = pd.read_excel(BytesIO(content), sheet_name=None, engine=engine)
        except ImportError as exc:
            raise ImportError(
                f"reading .{extension} attachments requires the `{engine}` engine; "
                f"run `pip install {engine}`"
            ) from exc

        parts = []
        for name, frame in sheets.items():
            parts.append(f"{name}:")
            parts.extend("\t".join(str(v) for v in row) for _, row in frame.iterrows())
        return "\n".join(parts).strip()

    def _extract_csv(self, content: bytes) -> str:
        try:
            import pandas as pd
        except ImportError as exc:
            raise ImportError(
                "`pandas` is required for CSV attachments; " \
                "run `pip install pandas`"
            ) from exc

        frame = pd.read_csv(BytesIO(content), low_memory=False)
        return "\n".join(", ".join(row.astype(str)) for _, row in frame.iterrows())

    def _extract_msg(self, content: bytes) -> str:
        try:
            import extract_msg
        except ImportError as exc:
            raise ImportError(
                "`extract-msg` is required for Outlook .msg attachments; "
                "run `pip install extract-msg`"
            ) from exc

        with extract_msg.Message(BytesIO(content)) as msg:
            return (
                f"Subject: {msg.subject}\nFrom: {msg.sender}\nTo: {msg.to}\n"
                f"CC: {msg.cc}\n\n{msg.body}"
            )

    def _extract_html(self, content: bytes) -> str:
        try:
            from bs4 import BeautifulSoup
        except ImportError as exc:
            raise ImportError(
                "`beautifulsoup4` is required for HTML attachments; "
                "run `pip install beautifulsoup4`"
            ) from exc

        return BeautifulSoup(content, "html.parser").get_text(separator=" ", strip=True)

    def _extract_text(self, content: bytes) -> str:
        return content.decode("utf-8", errors="replace")