from __future__ import annotations

from collections.abc import Iterable
from io import BytesIO
from typing import Any, ClassVar

from harborrag_core.domain.element import DocumentElement
from harborrag_core.domain.parser import ParsedDocument, ParseInput

from .base import BaseParser
from .exceptions import ParseError
from .parser_logging import get_parser_logger, input_label, parser_log_extra
from .utils import compact_text, guard_input_size, wrap_parse_errors

parser_logger = get_parser_logger("pptx")


class PptxParser(BaseParser[ParseInput, ParsedDocument]):
    """Extract slide text, tables, and grouped-shape text from PowerPoint files."""

    parser_name: ClassVar[str] = "pptx"
    parser_engine: ClassVar[str] = "python-pptx"
    suffixes: ClassVar[frozenset[str]] = frozenset({"pptx", "pptm"})
    content_types: ClassVar[frozenset[str]] = frozenset(
        {
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint.presentation.macroenabled.12",
        }
    )

    def parse(self, input: ParseInput) -> ParsedDocument:
        """Walk slides and return one document element per slide with text."""

        parse_input = self.coerce_input(input)
        try:
            from pptx import Presentation
        except ImportError as exc:
            parser_logger.error(
                "PPTX parser dependency `python-pptx` is missing",
                extra=parser_log_extra(
                    input=parse_input,
                    parser_name=self.parser_name,
                    parser_engine=self.parser_engine,
                ),
            )
            raise ParseError(
                "PPTX parsing requires `python-pptx`; install "
                "`harborrag-adapters[parsers]` or `pip install python-pptx`."
            ) from exc

        source_bytes = guard_input_size(parse_input.read_bytes())
        parser_logger.debug(
            "Extracting PPTX text from %s",
            input_label(parse_input),
            extra=parser_log_extra(
                input=parse_input,
                parser_name=self.parser_name,
                parser_engine=self.parser_engine,
                input_bytes=len(source_bytes),
            ),
        )
        sections: list[str] = []
        elements: list[DocumentElement] = []
        with wrap_parse_errors(self.parser_engine):
            presentation = Presentation(BytesIO(source_bytes))
            for slide_index, slide in enumerate(presentation.slides, start=1):
                slide_lines = list(self._shape_text(slide.shapes))
                slide_content = compact_text("\n".join(slide_lines))
                if not slide_content:
                    continue
                sections.append(f"Slide {slide_index}\n{slide_content}")
                elements.append(
                    DocumentElement(
                        id=f"pptx:slide:{slide_index}",
                        type="paragraph",
                        content=slide_content,
                        metadata={"slide": slide_index},
                    )
                )

        content = "\n\n".join(sections).strip()
        parser_logger.info(
            "Parsed PPTX %s slides=%d content_chars=%d elements=%d",
            input_label(parse_input),
            len(presentation.slides),
            len(content),
            len(elements),
            extra=parser_log_extra(
                input=parse_input,
                parser_name=self.parser_name,
                parser_engine=self.parser_engine,
                input_bytes=len(source_bytes),
                slides=len(presentation.slides),
                content_chars=len(content),
                elements=len(elements),
            ),
        )
        return ParsedDocument(
            content=content,
            elements=elements,
            parser_name=self.parser_name,
            parser_version=self.parser_version,
            metadata=self.metadata_for(
                parse_input,
                slide_count=len(presentation.slides),
            ),
        )

    @classmethod
    def _shape_text(cls, shapes: Iterable[Any], depth: int = 0) -> Iterable[str]:
        """Yield text from shapes, tables, and nested groups in display order."""

        for shape in shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if not text:
                        text = paragraph.text.strip()
                    if text:
                        yield text

            if getattr(shape, "has_table", False):
                table_lines = []
                for row in shape.table.rows:
                    table_lines.append(
                        "\t".join(cell.text.strip() for cell in row.cells).rstrip()
                    )
                table_text = "\n".join(line for line in table_lines if line.strip())
                if table_text:
                    yield table_text

            child_shapes = getattr(shape, "shapes", None)
            if child_shapes is not None and depth < 8:
                yield from cls._shape_text(child_shapes, depth + 1)
