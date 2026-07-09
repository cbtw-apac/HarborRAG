from __future__ import annotations

from io import BytesIO
from pathlib import Path
from tempfile import NamedTemporaryFile
from typing import Any, ClassVar, Iterable

from harborrag_core.domain.element import DocumentElement
from harborrag_core.domain.parser import ParsedDocument, ParseInput

from .base import BaseParser
from .exceptions import ParseError
from .parser_logging import get_parser_logger, input_label, parser_log_extra
from .utils import compact_text, guard_input_size, wrap_parse_errors


parser_logger = get_parser_logger("office")


class DocxParser(BaseParser[ParseInput, ParsedDocument]):
    """Extract text from Word `.docx` files with docx2txt."""

    parser_name: ClassVar[str] = "docx"
    parser_engine: ClassVar[str] = "docx2txt"
    suffixes: ClassVar[frozenset[str]] = frozenset({"docx"})
    content_types: ClassVar[frozenset[str]] = frozenset(
        {"application/vnd.openxmlformats-officedocument.wordprocessingml.document"}
    )

    def parse(self, input: ParseInput) -> ParsedDocument:
        """Materialize DOCX bytes to a temp path because docx2txt expects a file."""

        parse_input = self.coerce_input(input)
        try:
            import docx2txt
        except ImportError as exc:
            parser_logger.error(
                "DOCX parser dependency `docx2txt` is missing",
                extra=parser_log_extra(
                    input=parse_input,
                    parser_name=self.parser_name,
                    parser_engine=self.parser_engine,
                ),
            )
            raise ParseError(
                "DOCX parsing requires `docx2txt`; install `harborrag-adapters[parsers]` "
                "or `pip install docx2txt`."
            ) from exc

        tmp_path: Path | None = None
        try:
            parser_logger.debug(
                "Extracting DOCX text from %s",
                input_label(parse_input),
                extra=parser_log_extra(
                    input=parse_input,
                    parser_name=self.parser_name,
                    parser_engine=self.parser_engine,
                ),
            )
            with wrap_parse_errors(self.parser_engine):
                with NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
                    tmp.write(guard_input_size(parse_input.read_bytes()))
                    tmp_path = Path(tmp.name)
                content = compact_text(docx2txt.process(str(tmp_path)) or "")
        finally:
            if tmp_path is not None:
                tmp_path.unlink(missing_ok=True)

        elements = [
            DocumentElement(
                id="docx:0",
                type="paragraph",
                content=content,
                metadata={"filename": parse_input.filename},
            )
        ] if content else []
        return ParsedDocument(
            content=content,
            elements=elements,
            parser_name=self.parser_name,
            parser_version=self.parser_version,
            metadata=self.metadata_for(parse_input),
        )


class PptxParser(BaseParser[ParseInput, ParsedDocument]):
    """Extract slide text, tables, and grouped-shape text from PowerPoint files."""

    parser_name: ClassVar[str] = "pptx"
    parser_engine: ClassVar[str] = "python-pptx"
    suffixes: ClassVar[frozenset[str]] = frozenset({"pptx", "pptm"})
    content_types: ClassVar[frozenset[str]] = frozenset(
        {
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint.presentation.macroenabled.12",
        }
    )

    def parse(self, input: ParseInput) -> ParsedDocument:
        """Walk slides and return one document element per slide with text."""

        parse_input = self.coerce_input(input)
        try:
            from pptx import Presentation
        except ImportError as exc:
            parser_logger.error(
                "PPTX parser dependency `python-pptx` is missing",
                extra=parser_log_extra(
                    input=parse_input,
                    parser_name=self.parser_name,
                    parser_engine=self.parser_engine,
                ),
            )
            raise ParseError(
                "PPTX parsing requires `python-pptx`; install "
                "`harborrag-adapters[parsers]` or `pip install python-pptx`."
            ) from exc

        parser_logger.debug(
            "Extracting PPTX text from %s",
            input_label(parse_input),
            extra=parser_log_extra(
                input=parse_input,
                parser_name=self.parser_name,
                parser_engine=self.parser_engine,
            ),
        )
        sections: list[str] = []
        elements: list[DocumentElement] = []
        with wrap_parse_errors(self.parser_engine):
            presentation = Presentation(BytesIO(guard_input_size(parse_input.read_bytes())))
            for slide_index, slide in enumerate(presentation.slides, start=1):
                slide_lines = list(self._shape_text(slide.shapes))
                slide_content = compact_text("\n".join(slide_lines))
                if not slide_content:
                    continue
                sections.append(f"Slide {slide_index}\n{slide_content}")
                elements.append(
                    DocumentElement(
                        id=f"pptx:slide:{slide_index}",
                        type="paragraph",
                        content=slide_content,
                        metadata={"slide": slide_index},
                    )
                )

        content = "\n\n".join(sections).strip()
        return ParsedDocument(
            content=content,
            elements=elements,
            parser_name=self.parser_name,
            parser_version=self.parser_version,
            metadata=self.metadata_for(parse_input, slide_count=len(presentation.slides)),
        )

    @classmethod
    def _shape_text(cls, shapes: Iterable[Any], depth: int = 0) -> Iterable[str]:
        """Yield text from shapes, tables, and nested groups in display order."""

        for shape in shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if not text:
                        text = paragraph.text.strip()
                    if text:
                        yield text

            if getattr(shape, "has_table", False):
                table_lines = []
                for row in shape.table.rows:
                    table_lines.append(
                        "\t".join(cell.text.strip() for cell in row.cells).rstrip()
                    )
                table_text = "\n".join(line for line in table_lines if line.strip())
                if table_text:
                    yield table_text

            child_shapes = getattr(shape, "shapes", None)
            if child_shapes is not None and depth < 8:
                yield from cls._shape_text(child_shapes, depth + 1)


class ExcelParser(BaseParser[ParseInput, ParsedDocument]):
    """Extract workbook sheets as tab-separated table text."""

    parser_name: ClassVar[str] = "excel"
    parser_engine: ClassVar[str] = "openpyxl/xlrd"
    suffixes: ClassVar[frozenset[str]] = frozenset(
        {"xls", "xlsx", "xlsm", "xltx", "xltm"}
    )
    content_types: ClassVar[frozenset[str]] = frozenset(
        {
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
            "application/vnd.ms-excel.sheet.macroenabled.12",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.template",
            "application/vnd.ms-excel.template.macroenabled.12",
        }
    )

    def parse(self, input: ParseInput) -> ParsedDocument:
        """Route legacy `.xls` through xlrd and OpenXML workbooks through openpyxl."""

        parse_input = self.coerce_input(input)
        if parse_input.suffix == ".xls":
            content, elements, sheet_names = self._parse_xls(parse_input)
        else:
            content, elements, sheet_names = self._parse_openxml(parse_input)

        return ParsedDocument(
            content=content,
            elements=elements,
            parser_name=self.parser_name,
            parser_version=self.parser_version,
            metadata=self.metadata_for(parse_input, sheets=sheet_names),
        )

    def _parse_openxml(
        self, parse_input: ParseInput
    ) -> tuple[str, list[DocumentElement], list[str]]:
        """Read `.xlsx`-style workbooks in read-only, data-only mode."""

        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            parser_logger.error(
                "Excel parser dependency `openpyxl` is missing",
                extra=parser_log_extra(
                    input=parse_input,
                    parser_name=self.parser_name,
                    parser_engine=self.parser_engine,
                ),
            )
            raise ParseError(
                "Excel parsing requires `openpyxl`; install `harborrag-adapters[parsers]` "
                "or `pip install openpyxl`."
            ) from exc

        parser_logger.debug(
            "Extracting OpenXML Excel text from %s",
            input_label(parse_input),
            extra=parser_log_extra(
                input=parse_input,
                parser_name=self.parser_name,
                parser_engine=self.parser_engine,
            ),
        )
        with wrap_parse_errors(self.parser_engine):
            workbook = load_workbook(
                BytesIO(guard_input_size(parse_input.read_bytes())),
                read_only=True,
                data_only=True,
                keep_links=False,
            )
        sheet_names = workbook.sheetnames
        try:
            sections: list[str] = []
            elements: list[DocumentElement] = []
            for sheet in workbook.worksheets:
                rows = [
                    "\t".join(self._cell_to_text(value) for value in row).rstrip()
                    for row in sheet.iter_rows(values_only=True)
                ]
                rows = [row for row in rows if row.strip()]
                if not rows:
                    continue
                sheet_text = "\n".join(rows)
                sections.append(f"Sheet: {sheet.title}\n{sheet_text}")
                elements.append(
                    DocumentElement(
                        id=f"excel:sheet:{sheet.title}",
                        type="table",
                        content=sheet_text,
                        metadata={"sheet": sheet.title},
                    )
                )
        finally:
            workbook.close()

        return "\n\n".join(sections).strip(), elements, sheet_names

    def _parse_xls(
        self, parse_input: ParseInput
    ) -> tuple[str, list[DocumentElement], list[str]]:
        """Read legacy binary `.xls` workbooks with xlrd."""

        try:
            import xlrd
        except ImportError as exc:
            parser_logger.error(
                "Excel parser dependency `xlrd` is missing",
                extra=parser_log_extra(
                    input=parse_input,
                    parser_name=self.parser_name,
                    parser_engine=self.parser_engine,
                ),
            )
            raise ParseError(
                "Legacy .xls parsing requires `xlrd`; install "
                "`harborrag-adapters[parsers]` or `pip install xlrd`."
            ) from exc

        parser_logger.debug(
            "Extracting legacy Excel text from %s",
            input_label(parse_input),
            extra=parser_log_extra(
                input=parse_input,
                parser_name=self.parser_name,
                parser_engine=self.parser_engine,
            ),
        )
        with wrap_parse_errors(self.parser_engine):
            workbook = xlrd.open_workbook(
                file_contents=guard_input_size(parse_input.read_bytes()),
                on_demand=True,
            )
        sheet_names = workbook.sheet_names()
        sections: list[str] = []
        elements: list[DocumentElement] = []
        try:
            # Load sheets one at a time (Book.sheets() would defeat on_demand and
            # load them all), unloading each after rendering to bound memory.
            for sheet_index in range(workbook.nsheets):
                sheet = workbook.sheet_by_index(sheet_index)
                rows = []
                for row_index in range(sheet.nrows):
                    row = "\t".join(
                        self._xls_cell_to_text(
                            sheet.cell(row_index, column_index),
                            workbook.datemode,
                            xlrd,
                        )
                        for column_index in range(sheet.ncols)
                    ).rstrip()
                    if row.strip():
                        rows.append(row)
                if not rows:
                    continue
                sheet_text = "\n".join(rows)
                sections.append(f"Sheet: {sheet.name}\n{sheet_text}")
                elements.append(
                    DocumentElement(
                        id=f"excel:sheet:{sheet.name}",
                        type="table",
                        content=sheet_text,
                        metadata={"sheet": sheet.name},
                    )
                )
                workbook.unload_sheet(sheet_index)
        finally:
            workbook.release_resources()

        return "\n\n".join(sections).strip(), elements, sheet_names

    @staticmethod
    def _cell_to_text(value: Any) -> str:
        """Convert openpyxl cell values into stable searchable text."""

        if value is None:
            return ""
        isoformat = getattr(value, "isoformat", None)
        if callable(isoformat):
            return isoformat()
        return str(value)

    @staticmethod
    def _xls_cell_to_text(cell: Any, datemode: int, xlrd: Any) -> str:
        """Convert xlrd cell values while preserving dates, booleans, and errors."""

        if cell.ctype in {xlrd.XL_CELL_EMPTY, xlrd.XL_CELL_BLANK}:
            return ""
        if cell.ctype == xlrd.XL_CELL_DATE:
            try:
                return xlrd.xldate_as_datetime(cell.value, datemode).isoformat()
            except (OverflowError, ValueError):
                return str(cell.value)
        if cell.ctype == xlrd.XL_CELL_NUMBER:
            number = float(cell.value)
            return str(int(number)) if number.is_integer() else str(number)
        if cell.ctype == xlrd.XL_CELL_BOOLEAN:
            return "TRUE" if cell.value else "FALSE"
        if cell.ctype == xlrd.XL_CELL_ERROR:
            return f"#ERROR:{cell.value}"
        return str(cell.value)
