from __future__ import annotations

import io
import struct
import zipfile
from collections.abc import Iterator
from dataclasses import dataclass, field
from html import escape
from typing import Any


def build_docx_bytes(text: str = "Hello Harbor") -> bytes:
    from docx import Document

    document = Document()
    document.add_paragraph(text)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def build_pptx_bytes(text: str = "Slide text") -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = text
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def build_xlsx_bytes(rows: list[list[Any]] | None = None) -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Sheet1"
    effective_rows = [["name", "role"], ["Ada", "engineer"]] if rows is None else rows
    for row in effective_rows:
        sheet.append(row)
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def build_epub_bytes(sections: list[str] | None = None) -> bytes:
    sections = ["Chapter one text", "Chapter two text"] if sections is None else sections
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr(
            "mimetype",
            "application/epub+zip",
            compress_type=zipfile.ZIP_STORED,
        )
        archive.writestr(
            "META-INF/container.xml",
            '<?xml version="1.0"?>'
            '<container version="1.0" '
            'xmlns="urn:oasis:names:tc:opendocument:xmlns:container">'
            '<rootfiles><rootfile full-path="OEBPS/content.opf" '
            'media-type="application/oebps-package+xml"/></rootfiles></container>',
        )
        manifest_items = []
        spine_items = []
        for index in range(1, len(sections) + 1):
            item_id = f"c{index}"
            href = f"ch{index}.xhtml"
            manifest_items.append(
                f'<item id="{item_id}" href="{href}" media-type="application/xhtml+xml"/>'
            )
            spine_items.append(f'<itemref idref="{item_id}"/>')
            archive.writestr(
                f"OEBPS/{href}",
                f"<html><body><p>{escape(sections[index - 1])}</p></body></html>",
            )
        archive.writestr(
            "OEBPS/content.opf",
            '<?xml version="1.0"?>'
            '<package xmlns="http://www.idpf.org/2007/opf" version="3.0" '
            'unique-identifier="id"><metadata/>'
            f"<manifest>{''.join(manifest_items)}</manifest>"
            f"<spine>{''.join(spine_items)}</spine></package>",
        )
    return buffer.getvalue()


def build_zip_bomb_bytes(member_size: int = 50 * 1024 * 1024, members: int = 4) -> bytes:
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        payload = b"\0" * member_size
        for index in range(members):
            archive.writestr(f"member{index}.txt", payload)
    return buffer.getvalue()


def build_zero_compressed_size_zip_bytes(
    filename: str = "bomb.bin", *, claimed_file_size: int = 10_000_000
) -> bytes:
    """Forge a ZIP whose central directory claims 0 compressed bytes but a
    large uncompressed size -- an effectively infinite ratio that a
    compression-ratio-only check would divide-by-zero around or skip.

    A real ``ZIP_STORED`` member is written first (so the archive is
    otherwise well-formed), then its local file header and central directory
    record are patched in place: compressed size -> 0, uncompressed size ->
    ``claimed_file_size``.
    """
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_STORED) as archive:
        archive.writestr(filename, b"x")
    data = bytearray(buffer.getvalue())

    with zipfile.ZipFile(io.BytesIO(bytes(data))) as archive:
        local_header_offset = archive.infolist()[0].header_offset

    # Local file header: compressed size at +18, uncompressed size at +22.
    struct.pack_into("<I", data, local_header_offset + 18, 0)
    struct.pack_into("<I", data, local_header_offset + 22, claimed_file_size)

    # Central directory record: compressed size at +20, uncompressed at +24.
    central_directory_offset = data.index(b"PK\x01\x02")
    struct.pack_into("<I", data, central_directory_offset + 20, 0)
    struct.pack_into("<I", data, central_directory_offset + 24, claimed_file_size)

    return bytes(data)


def build_understated_file_size_zip_bytes(
    real_content: bytes, *, claimed_file_size: int, filename: str = "member.bin"
) -> bytes:
    """Forge a ZIP whose central directory UNDER-reports uncompressed size.

    ``compress_size`` is left accurate (a real DEFLATE stream for
    ``real_content``), so metadata-only checks (total bytes, compression
    ratio) see a small, unremarkable ``claimed_file_size`` and pass -- only
    actually decompressing the member reveals it produces far more bytes
    than the forged metadata claims.
    """
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr(filename, real_content)
    data = bytearray(buffer.getvalue())

    with zipfile.ZipFile(io.BytesIO(bytes(data))) as archive:
        local_header_offset = archive.infolist()[0].header_offset

    # Local file header: uncompressed size at +22 (compressed size at +18 is
    # left untouched -- it must stay accurate for the real member to decode).
    struct.pack_into("<I", data, local_header_offset + 22, claimed_file_size)

    # Central directory record: uncompressed size at +24 (+20 untouched).
    central_directory_offset = data.index(b"PK\x01\x02")
    struct.pack_into("<I", data, central_directory_offset + 24, claimed_file_size)

    return bytes(data)


@dataclass
class FakeResponse:
    """Minimal stand-in for ``requests.Response`` used by client-level tests."""

    status_code: int = 200
    headers: dict[str, str] = field(default_factory=dict)
    _json: Any = None
    text: str = ""
    _chunks: list[bytes] = field(default_factory=list)
    closed: bool = False

    def json(self) -> Any:
        if self._json is None:
            raise ValueError("no json")
        return self._json

    def iter_content(self, chunk_size: int = 65536) -> Iterator[bytes]:
        _ = chunk_size
        yield from self._chunks

    def close(self) -> None:
        self.closed = True

    @property
    def content(self) -> bytes:
        return b"".join(self._chunks)


@dataclass
class FakeSession:
    """Records requests and replays a scripted queue of responses/exceptions."""

    responses: list[Any] = field(default_factory=list)
    calls: list[dict[str, Any]] = field(default_factory=list)
    headers: dict[str, str] = field(default_factory=dict)
    auth: Any = None

    def request(self, method: str, url: str, **kwargs: Any):
        self.calls.append({"method": method, "url": url, **kwargs})
        item = self.responses.pop(0)
        if isinstance(item, Exception):
            raise item
        return item

    def post(self, url: str, **kwargs: Any):
        """Support clients that call ``session.post`` directly (e.g. OAuth token exchange)."""
        return self.request("POST", url, **kwargs)
